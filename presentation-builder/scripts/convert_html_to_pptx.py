"""
Convert Stackd Logistics HTML presentation to PowerPoint PPTX format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# FirstMile brand colors
FM_BLUE = RGBColor(24, 42, 90)  # #182A5A
FM_GREEN = RGBColor(139, 215, 78)  # #8BD74E
FM_DARK = RGBColor(15, 24, 33)  # #0f1821
FM_LIGHT = RGBColor(224, 232, 240)  # #e0e8f0

def create_presentation():
    """Create the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(16)  # Widescreen 16:9
    prs.slide_height = Inches(9)

    # Slide 1: Title
    slide = create_title_slide(prs)

    # Slide 2: Your Shipping Profile
    slide = create_shipping_profile_slide(prs)

    # Slide 3: Network Expansion
    slide = create_network_expansion_slide(prs)

    # Slide 4: Xparcel Ground Rates
    slide = create_ground_rates_slide(prs)

    # Slide 5: Xparcel Expedited Rates
    slide = create_expedited_rates_slide(prs)

    # Slide 6: Zone Distribution
    slide = create_zone_distribution_slide(prs)

    # Slide 7: Service Levels
    slide = create_service_levels_slide(prs)

    # Slide 8: Savings Projection
    slide = create_savings_slide(prs)

    # Slide 9: Where FirstMile Shines
    slide = create_firstmile_shines_slide(prs)

    # Slide 10: Service Advantages
    slide = create_service_advantages_slide(prs)

    # Slide 11: Implementation
    slide = create_implementation_slide(prs)

    # Slide 12: Next Steps
    slide = create_next_steps_slide(prs)

    # Slide 13: Decision Point
    slide = create_decision_slide(prs)

    # Slide 14: Call to Action
    slide = create_cta_slide(prs)

    return prs

def set_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, top=0.5, color=FM_GREEN):
    """Add a title to the slide"""
    left = Inches(0.5)
    width = Inches(15)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = 'Maven Pro'

def add_text(slide, text, left, top, width, height, size=20, color=FM_LIGHT, bold=False):
    """Add text to the slide"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Maven Pro'

    return text_box

def create_title_slide(prs):
    """Slide 1: Title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    # Main title
    add_title(slide, "Stackd Logistics", 2.5, RGBColor(255, 255, 255))

    # Subtitle
    left = Inches(0.5)
    top = Inches(3.5)
    width = Inches(15)
    height = Inches(1)

    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FirstMile Xparcel Proposal"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = FM_GREEN
    p.font.name = 'Maven Pro'

    # Prepared for
    add_text(slide, "Prepared for: Landon Richards", 0.5, 5.5, 15, 0.5, 22, FM_LIGHT)
    add_text(slide, "8,957 Shipments Analyzed | 17,914 Packages/Month | ~896/Day",
             0.5, 6.2, 15, 0.5, 18, FM_GREEN)

    return slide

def create_shipping_profile_slide(prs):
    """Slide 2: Your Shipping Profile"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Your Current Shipping Reality", 0.5)

    # Metrics in a 2x2 grid
    metrics = [
        ("17,914", "Packages/Month"),
        ("~896", "Average Daily Volume"),
        ("$4.96", "Average Cost/Label"),
        ("89.9%", "Packages Under 1 lb")
    ]

    positions = [(1, 1.8), (8.5, 1.8), (1, 3.5), (8.5, 3.5)]

    for (value, label), (left, top) in zip(metrics, positions):
        # Value
        add_text(slide, value, left, top, 6, 0.6, 42, FM_GREEN, bold=True)
        # Label
        add_text(slide, label, left, top + 0.7, 6, 0.4, 16, FM_LIGHT)

    # Weight breakdown
    add_text(slide, "Weight Distribution Breakdown", 0.5, 5.5, 15, 0.5, 28, RGBColor(255, 255, 255), bold=True)

    bullets = [
        "• 55% of shipments are 0-4 ounces",
        "• 27% are 4-8 ounces",
        "• 10% are 8 ounces to 1 pound",
        "• Only 10.1% are over 1 pound"
    ]

    y = 6.2
    for bullet in bullets:
        add_text(slide, bullet, 0.5, y, 15, 0.4, 20, FM_LIGHT)
        y += 0.5

    return slide

def create_network_expansion_slide(prs):
    """Slide 3: Network Expansion"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "The Expansion of FirstMile's Xparcel Networks", 0.5)
    add_text(slide, "From 12 injection points to 47 metro markets nationwide",
             0.5, 1.3, 15, 0.4, 20, FM_GREEN)

    # OLD Network box
    from pptx.util import Inches
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(6.5)
    height = Inches(4)

    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 100, 100)
    shape.fill.transparency = 0.9
    shape.line.color.rgb = RGBColor(255, 107, 107)
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "OLD Network\n\n• 12 injection points\n• Limited metro coverage\n• Longer zone distances\n• Higher transportation costs\n• Fewer routing options"
    p.font.size = Pt(18)
    p.font.color.rgb = FM_LIGHT
    p.font.name = 'Maven Pro'

    # NEW Network box
    left = Inches(8.5)
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = FM_GREEN
    shape.fill.transparency = 0.85
    shape.line.color.rgb = FM_GREEN
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NEW Network (2025)\n\n• 47 metro injection points\n• Covers ~70% of U.S. population\n• Zone-skipping advantages\n• Reduced average spend per package\n• Intelligent routing to closest hub"
    p.font.size = Pt(18)
    p.font.color.rgb = FM_LIGHT
    p.font.name = 'Maven Pro'

    return slide

def create_ground_rates_slide(prs):
    """Slide 4: Xparcel Ground Rates"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Xparcel Ground Rates - Your Primary Service", 0.5)
    add_text(slide, "3-8 Day Delivery | 96% of Your Volume", 0.5, 1.3, 15, 0.4, 24, FM_GREEN)

    # Table data
    rows = 6
    cols = 5
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(12)
    height = Inches(3.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header row
    headers = ["Weight", "Zone 4", "Zone 5", "Zone 6", "Zone 7"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = FM_BLUE
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data rows
    data = [
        ["1-4 oz", "$4.02", "$4.07", "$4.16", "$4.23"],
        ["5-8 oz", "$4.31", "$4.35", "$4.40", "$4.40"],
        ["9-12 oz", "$4.60", "$4.72", "$4.91", "$5.04"],
        ["13-15.99 oz", "$5.19", "$5.41", "$5.70", "$5.92"],
        ["1 lb (16 oz)", "$5.84", "$6.11", "$6.30", "$6.65"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 40, 70)
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].font.color.rgb = FM_LIGHT

    # Add comparison note
    add_text(slide, "Current Average: $4.96/package", 2, 6.5, 5, 0.4, 18, FM_LIGHT, bold=True)
    add_text(slide, "Xparcel Ground: $4.20-$4.60/package", 9, 6.5, 5, 0.4, 18, FM_GREEN, bold=True)

    return slide

def create_expedited_rates_slide(prs):
    """Slide 5: Xparcel Expedited Rates"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Xparcel Expedited - Faster Ground Option", 0.5)
    add_text(slide, "2-5 Day Delivery | 4% of Your Volume", 0.5, 1.3, 15, 0.4, 24, FM_GREEN)

    add_text(slide, "Select Network (High-Density Markets)", 0.5, 2, 15, 0.4, 20, RGBColor(255, 255, 255), bold=True)

    # Select Network table
    rows = 5
    cols = 5
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(13)
    height = Inches(2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header
    headers = ["Weight", "Zone 4", "Zone 5", "Zone 6", "Zone 7"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = FM_BLUE
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data
    data = [
        ["1 lb", "$3.29", "$3.34", "$3.35", "$3.38"],
        ["2 lbs", "$3.35", "$3.38", "$3.40", "$3.45"],
        ["3 lbs", "$3.41", "$3.45", "$3.56", "$3.35"],
        ["5 lbs", "$3.48", "$3.52", "$3.56", "$3.65"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 40, 70)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = FM_LIGHT

    # Note
    add_text(slide, "30-40% below traditional carrier Express services for equivalent transit times",
             1.5, 7.5, 13, 0.6, 20, FM_GREEN, bold=True)

    return slide

def create_zone_distribution_slide(prs):
    """Slide 6: Zone Distribution"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "The Geographic Challenge", 0.5)
    add_text(slide, "Your zone distribution reveals a significant opportunity:",
             0.5, 1.3, 15, 0.4, 20, FM_LIGHT)

    # Zone bars
    zones = [
        ("Zones 1-3 (Local)", "11%", 11),
        ("Zones 4-5 (Regional)", "40.5%", 40.5),
        ("Zones 6-7 (Cross-Country)", "47.8%", 47.8),
        ("Zone 8 (Far Reach)", "1.2%", 1.2)
    ]

    y = 3
    for label, percentage, width_pct in zones:
        # Label
        add_text(slide, label, 1, y, 4, 0.4, 18, RGBColor(255, 255, 255), bold=True)

        # Bar background
        from pptx.util import Inches
        bar_bg = slide.shapes.add_shape(1, Inches(5.5), Inches(y), Inches(9), Inches(0.4))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = RGBColor(50, 60, 90)
        bar_bg.line.color.rgb = RGBColor(50, 60, 90)

        # Bar fill
        bar_width = 9 * (width_pct / 100)
        bar = slide.shapes.add_shape(1, Inches(5.5), Inches(y), Inches(bar_width), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = FM_GREEN
        bar.line.color.rgb = FM_GREEN

        # Percentage text
        add_text(slide, percentage, 5.5 + bar_width + 0.2, y, 2, 0.4, 16, FM_GREEN, bold=True)

        y += 1

    add_text(slide, "89% of your volume ships to zones 4-8 — where traditional carriers charge premium zone penalties.",
             1, 7.5, 14, 0.8, 20, FM_GREEN)

    return slide

def create_service_levels_slide(prs):
    """Slide 7: Service Levels"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Xparcel Service Levels", 0.5)
    add_text(slide, "Three service tiers designed specifically for eCommerce shippers:",
             0.5, 1.3, 15, 0.4, 20, FM_LIGHT)

    # Table
    rows = 4
    cols = 4
    left = Inches(1.5)
    top = Inches(2.5)
    width = Inches(13)
    height = Inches(3)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header
    headers = ["Service Level", "Transit Time", "Best For", "Your Volume Fit"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = FM_BLUE
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data
    data = [
        ["Xparcel Ground", "3-8 days", "Economy shipments, under 1 lb", "96% of your volume"],
        ["Xparcel Expedited", "2-5 days", "Faster ground, 1-20 lbs", "4% of your volume"],
        ["Xparcel Priority", "1-3 days", "Premium speed with guarantee", "Available for urgent needs"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 40, 70)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = FM_LIGHT

    return slide

def create_savings_slide(prs):
    """Slide 8: Savings Projection"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Your Conservative Savings Projection", 0.5)

    # Table
    rows = 8
    cols = 3
    left = Inches(2)
    top = Inches(2)
    width = Inches(12)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header
    headers = ["Metric", "Current State", "FirstMile Xparcel"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = FM_BLUE
        cell.text_frame.paragraphs[0].font.size = Pt(18)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data
    data = [
        ["Monthly Packages", "17,914", "17,914"],
        ["Daily Volume (5 days/wk)", "~896", "~896"],
        ["Average Cost/Package", "$4.96", "$4.20 - $4.60"],
        ["Monthly Spend", "$88,853", "$75,200 - $82,400"],
        ["Monthly Savings", "—", "$6,500 - $13,700"],
        ["Annual Savings", "—", "$78,000 - $164,000"],
        ["Cost Reduction %", "—", "7% - 15%"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 40, 70)
            cell.text_frame.paragraphs[0].font.size = Pt(16)

            # Highlight savings rows in green
            if row_idx >= 5:
                cell.text_frame.paragraphs[0].font.color.rgb = FM_GREEN
            else:
                cell.text_frame.paragraphs[0].font.color.rgb = FM_LIGHT

    add_text(slide, "Conservative projections based on your actual shipping data",
             2, 7, 12, 0.5, 18, FM_GREEN)

    return slide

def create_firstmile_shines_slide(prs):
    """Slide 9: Where FirstMile Shines"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Where FirstMile Shines for Stackd", 0.5)

    # 4 boxes in 2x2 grid
    boxes = [
        ("Perfect Weight Profile", "89.9% under 1 lb is our sweet spot. Traditional carriers overprice lightweight packages."),
        ("Zone Distribution Advantage", "88.3% of volume in zones 4-7 is ideal for our expanded network."),
        ("Service Mix Simplicity", "96% Ground volume means we optimize the vast majority without complexity."),
        ("Strong Volume", "17,914 packages/month provides stable operations and strong negotiating position.")
    ]

    positions = [(1, 2.5), (8.5, 2.5), (1, 5.5), (8.5, 5.5)]

    for (title, desc), (left, top) in zip(boxes, positions):
        # Box background
        from pptx.util import Inches
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(6), Inches(2.3))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 40, 70)
        box.line.color.rgb = FM_GREEN
        box.line.width = Pt(2)

        # Title
        add_text(slide, title, left + 0.3, top + 0.2, 5.4, 0.5, 20, FM_GREEN, bold=True)

        # Description
        add_text(slide, desc, left + 0.3, top + 0.8, 5.4, 1.3, 16, FM_LIGHT)

    return slide

def create_service_advantages_slide(prs):
    """Slide 10: Service Advantages"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Beyond Cost: The Service Advantages", 0.5)
    add_text(slide, "Cost savings are important, but FirstMile provides operational benefits:",
             0.5, 1.3, 15, 0.4, 18, FM_LIGHT)

    bullets = [
        "• Audit Queue catches mis-rated labels before they hit your invoice",
        "• Single Support Thread — one FirstMile team handles claims, returns, exceptions",
        "• Returns Portal with QR-code-based returns requiring no printing",
        "• Claims Recovery at 94% success rate within 14 days"
    ]

    y = 2.5
    for bullet in bullets:
        add_text(slide, bullet, 1, y, 14, 0.6, 20, FM_LIGHT)
        y += 1

    add_text(slide, "For a 3PL like Stackd: Every dollar saved on shipping strengthens your service offering.",
             1, 7, 14, 0.8, 22, FM_GREEN, bold=True)

    return slide

def create_implementation_slide(prs):
    """Slide 11: Implementation"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "Implementation: Simpler Than You Think", 0.5)
    add_text(slide, "You're using ShipHero. FirstMile integrates directly through their multi-carrier API.",
             0.5, 1.3, 15, 0.4, 18, FM_LIGHT)

    # Timeline table
    rows = 6
    cols = 2
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(12)
    height = Inches(4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header
    headers = ["Week", "Milestone"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = FM_BLUE
        cell.text_frame.paragraphs[0].font.size = Pt(18)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Data
    data = [
        ["Week 1", "Contract signature, tech kickoff with ShipHero"],
        ["Weeks 2-3", "API integration dev and testing in sandbox"],
        ["Week 4", "Pilot launch with 20-30% volume"],
        ["Weeks 5-12", "Ramp to full volume with DHL parallel"],
        ["Week 13+", "Ongoing optimization, quarterly reviews"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 40, 70)
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].font.color.rgb = FM_LIGHT

    add_text(slide, "90-Day Pilot: Keep DHL running in parallel. Make the full switch only when you're confident.",
             2, 7, 12, 0.6, 18, FM_GREEN, bold=True)

    return slide

def create_next_steps_slide(prs):
    """Slide 12: Next Steps"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "What Happens Next", 0.5)
    add_text(slide, "If this makes sense based on the data — and it should — here's the path forward:",
             0.5, 1.3, 15, 0.4, 20, FM_LIGHT)

    # 4 boxes
    steps = [
        ("Today", ["Verbal commitment to 90-day pilot", "Provide 10-20 sample DHL invoices", "Answer key questions"]),
        ("This Week", ["Formal proposal with detailed rate cards", "Implementation timeline", "Contract terms review"]),
        ("Week 1", ["Contract signature", "Tech kickoff with ShipHero", "FirstMile account setup"]),
        ("Weeks 2-4", ["Integration dev and testing", "Pilot launch with 20-30% volume", "Daily monitoring"])
    ]

    positions = [(1, 2.5), (8.5, 2.5), (1, 5.3), (8.5, 5.3)]

    for (title, items), (left, top) in zip(steps, positions):
        # Box
        from pptx.util import Inches
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(6), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(30, 40, 70)
        box.line.color.rgb = FM_GREEN
        box.line.width = Pt(2)

        # Title
        add_text(slide, title, left + 0.3, top + 0.2, 5.4, 0.4, 20, FM_GREEN, bold=True)

        # Items
        y_offset = 0.7
        for item in items:
            add_text(slide, f"• {item}", left + 0.3, top + y_offset, 5.4, 0.4, 14, FM_LIGHT)
            y_offset += 0.45

    return slide

def create_decision_slide(prs):
    """Slide 13: The Decision Point"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    add_title(slide, "The Decision Point", 0.5)
    add_text(slide, "Landon, you have a decision to make:", 0.5, 1.3, 15, 0.4, 22, FM_LIGHT)

    # Two comparison boxes
    # Continue Current Path
    from pptx.util import Inches
    left_box = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(6.5), Inches(4))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(255, 100, 100)
    left_box.fill.transparency = 0.85
    left_box.line.color.rgb = RGBColor(255, 107, 107)
    left_box.line.width = Pt(3)

    add_text(slide, "Continue Current Path", 1.3, 2.8, 6, 0.5, 24, RGBColor(255, 107, 107), bold=True)

    current_items = [
        "• Pay $88,853 monthly",
        "• Traditional carrier structure",
        "• Zone penalties on 89% of volume",
        "• $1,066,236 annual shipping costs"
    ]

    y = 3.6
    for item in current_items:
        add_text(slide, item, 1.3, y, 6, 0.4, 16, RGBColor(255, 200, 200))
        y += 0.6

    # Optimize with FirstMile
    right_box = slide.shapes.add_shape(1, Inches(8.5), Inches(2.5), Inches(6.5), Inches(4))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = FM_GREEN
    right_box.fill.transparency = 0.85
    right_box.line.color.rgb = FM_GREEN
    right_box.line.width = Pt(3)

    add_text(slide, "Optimize with FirstMile", 8.8, 2.8, 6, 0.5, 24, FM_GREEN, bold=True)

    firstmile_items = [
        "• Pay $75,200-$82,400 monthly",
        "• eCommerce-optimized network",
        "• Zone-skipping on far destinations",
        "• $78K-$164K savings annually"
    ]

    y = 3.6
    for item in firstmile_items:
        add_text(slide, item, 8.8, y, 6, 0.4, 16, FM_LIGHT)
        y += 0.6

    add_text(slide, "Every month you wait costs you $6,500 to $13,700.",
             1, 7.2, 14, 0.6, 24, FM_GREEN, bold=True)

    return slide

def create_cta_slide(prs):
    """Slide 14: Call to Action"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide, FM_BLUE)

    # Centered content
    add_text(slide, "Ready to Optimize Your Shipping?", 3, 2.5, 10, 1, 48, RGBColor(255, 255, 255), bold=True)
    add_text(slide, "Let's start your 90-day pilot and prove the savings.", 3, 3.8, 10, 0.6, 28, FM_GREEN)

    # Contact info
    add_text(slide, "Brett Walker", 3, 5.5, 10, 0.4, 22, RGBColor(255, 255, 255))
    add_text(slide, "FirstMile Sales Executive", 3, 6, 10, 0.4, 18, FM_GREEN)
    add_text(slide, "brett@firstmile.com", 3, 6.6, 10, 0.4, 18, FM_LIGHT)

    return slide

def main():
    """Main function to create and save the presentation"""
    print("Creating PowerPoint presentation...")
    prs = create_presentation()

    # Save the presentation
    output_path = r"C:\Users\BrettWalker\FirstMile_Deals\presentation-builder\Finished Presentations\stackd-logistics-xparcel-presentation.pptx"
    prs.save(output_path)
    print(f"SUCCESS: Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
